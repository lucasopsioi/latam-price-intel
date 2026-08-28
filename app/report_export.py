# -*- coding: utf-8 -*-
"""周报导出：Word / PPT / PDF。

★★ 三种格式共用同一份**结构化解析**，而不是各写一遍 markdown 转换。
   报告的 markdown 是我们自己生成的（格式可控），所以这里的解析器只认
   我们会产出的那几种块：标题 / 段落 / 表格 / 图占位 / 引言 / 分隔线。
   不做通用 markdown 解析 —— 那是另一个项目的工作量，而且用不上。

★ 图在网页上是 ECharts，导出时必须服务端重画。这里用 matplotlib 复刻
   **同样的语义**（compare→哑铃图、share→100% 堆叠），而不是随手换个图形：
   同一份报告在网页和 Word 里长得不一样，读者会以为是两份数据。

★ 中文字体必须显式指定。matplotlib 默认字体没有中文字形，
   不指定会得到满屏方块 —— 而且**不报错**，导出的文件看起来"成功了"。
"""
from __future__ import annotations

import io
import re
from datetime import datetime

# ── matplotlib 必须在导入 pyplot 之前选无界面后端 ──
# 服务进程没有窗口系统，用默认后端会在某些环境下尝试建窗口而卡住。
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt                      # noqa: E402
from matplotlib import font_manager                  # noqa: E402

_CJK = ["Microsoft YaHei", "SimHei", "Noto Sans CJK SC", "Source Han Sans SC",
        "SimSun", "DengXian"]
_AVAIL = {f.name for f in font_manager.fontManager.ttflist}
_FONT = next((f for f in _CJK if f in _AVAIL), None)
if _FONT:
    plt.rcParams["font.sans-serif"] = [_FONT]
plt.rcParams["axes.unicode_minus"] = False           # 否则负号显示成方块

# 与网页版一致的配色（charts.js 的 accent / red / green）
C_ACCENT, C_DOWN, C_UP, C_GREY = "#0A84FF", "#34C759", "#FF3B30", "#8E8E93"


# ══════════════════════════ markdown → 结构块 ══════════════════════════

_MD_LINK = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")


def _strip_links(text: str) -> str:
    """[来源名](url) → 来源名。

    ★ 导出物里必须剥掉链接：Google News 的跳转 URL 几百个字符，
      塞进 Word/PPT 表格会把整列挤爆。网页端保留可点链接，
      导出端只留来源名 —— 要核对原文回网页版点。
    """
    return _MD_LINK.sub(r"\1", text or "")


def parse_blocks(md: str) -> list[dict]:
    """把报告 markdown 切成结构块。只认我们自己会产出的那几种。"""
    # ★ 抓来的商品标题/评论里会混控制字符（实测 \x0b 之类）——
    #   python-docx 的 lxml 直接 ValueError 拒收，PDF/PPT 则静默吞掉。
    #   在入口统一清洗，三种导出同一条防线。
    md = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", md or "")
    blocks: list[dict] = []
    lines = md.splitlines()
    i, n = 0, len(lines)
    para: list[str] = []

    def flush():
        if para:
            blocks.append({"t": "p", "text": _strip_links(" ".join(para).strip())})
            para.clear()

    while i < n:
        raw = lines[i]
        line = raw.strip()
        if not line:
            flush(); i += 1; continue
        m = re.match(r"^(#{1,4})\s+(.*)$", line)
        if m:
            flush()
            blocks.append({"t": "h", "level": len(m.group(1)), "text": m.group(2).strip()})
            i += 1; continue
        m = re.match(r"^!\[chart:([\w-]+)\]$", line)
        if m:
            flush()
            blocks.append({"t": "chart", "el": m.group(1)})
            i += 1; continue
        if line.startswith(">"):
            flush()
            blocks.append({"t": "quote", "text": line.lstrip("> ").strip()})
            i += 1; continue
        if line.startswith("---"):
            flush()
            blocks.append({"t": "hr"})
            i += 1; continue
        if line.startswith("|"):
            flush()
            rows = []
            while i < n and lines[i].strip().startswith("|"):
                cells = [_strip_links(c.strip()) for c in lines[i].strip().strip("|").split("|")]
                rows.append(cells)
                i += 1
            # 第 2 行是 |---| 分隔线，丢掉
            if len(rows) >= 2 and all(set(c) <= set("-: ") for c in rows[1]):
                rows = [rows[0]] + rows[2:]
            if rows:
                blocks.append({"t": "table", "head": rows[0], "rows": rows[1:]})
            continue
        para.append(re.sub(r"\*\*(.+?)\*\*", r"\1", line))
        i += 1
    flush()
    return blocks


# ══════════════════════════ 图：matplotlib 复刻 ══════════════════════════

def _fig_png(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=170, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    return buf.getvalue()


def _chart_deviation(opt: dict, xlab: str = "") -> bytes:
    """发散条：相对 0 的偏离（变动幅度 %）。与网页 deviation 同语义。

    ★ 用它而不是绝对价哑铃图，是因为报告覆盖六国六币种 ——
      绝对价同轴会被最大面值的币种压平（实测哥伦比亚 210 万压掉其余四国）。
    """
    rows = (opt.get("rows") or [])[:12]
    fig, ax = plt.subplots(figsize=(9, max(2.2, 0.42 * len(rows) + 1)))
    labels = [r.get("label", "") for r in reversed(rows)]
    vals = [float(r.get("v") or 0) for r in reversed(rows)]
    up_bad = opt.get("upIsBad", True)
    colors_ = [(C_UP if (v > 0) == bool(up_bad) else C_DOWN) for v in vals]
    ax.barh(labels, vals, height=0.55, color=colors_)
    ax.axvline(0, color="#3A3A3C", lw=1.0)
    for k, v in enumerate(vals):
        ax.text(v + (1.2 if v >= 0 else -1.2), k, f"{v:+.1f}%",
                va="center", ha="left" if v >= 0 else "right", fontsize=8.5)
    pad = max(abs(min(vals, default=0)), abs(max(vals, default=0))) * 0.28 + 2
    ax.set_xlim(min(vals, default=0) - pad, max(vals, default=0) + pad)
    ax.set_xlabel(xlab or "变动幅度 %", fontsize=9)
    ax.grid(axis="x", color="#E5E5EA", lw=0.7)
    ax.set_axisbelow(True)
    for sp in ("top", "right", "left"):
        ax.spines[sp].set_visible(False)
    ax.tick_params(labelsize=9)
    return _fig_png(fig)


def _chart_compare(opt: dict, xlab: str = "") -> bytes:
    """哑铃图：前 → 后。与网页 compare 同语义。"""
    rows = (opt.get("rows") or [])[:12]
    fig, ax = plt.subplots(figsize=(9, max(2.2, 0.42 * len(rows) + 1)))
    for k, r in enumerate(reversed(rows)):
        a, b = r.get("from"), r.get("to")
        if a is None or b is None:
            continue
        down = b < a
        ax.plot([a, b], [k, k], color=C_GREY, lw=1.4, zorder=1)
        ax.scatter([a], [k], s=42, facecolors="white",
                   edgecolors=C_GREY, zorder=2, linewidths=1.4)
        ax.scatter([b], [k], s=54, color=(C_DOWN if down else C_UP), zorder=3)
    ax.set_yticks(range(len(rows)))
    ax.set_yticklabels([r.get("label", "") for r in reversed(rows)], fontsize=9)
    ax.set_xlabel(xlab or "价格", fontsize=9)
    ax.grid(axis="x", color="#E5E5EA", lw=0.7)
    ax.set_axisbelow(True)
    for s in ("top", "right", "left"):
        ax.spines[s].set_visible(False)
    ax.tick_params(axis="x", labelsize=8)
    return _fig_png(fig)


def _chart_change(opt: dict, ylab: str = "") -> bytes:
    """多品牌折线（周度 ASP 等时间序列）。与网页 change 同语义。"""
    xs = opt.get("xs") or []
    series = opt.get("series") or []
    if not xs or not series:
        return b""
    fig, ax = plt.subplots(figsize=(8.6, 3.4))
    palette = ["#2a78d6", "#eb6834", "#4a3aa7", "#1baf7a", "#eda100",
               "#e87ba4", "#008300", "#8e8e93"]
    HW = "#C7000B"
    ends = []                    # (y, 文本, 颜色) —— 末端价格标注，防重叠用
    for i, sr in enumerate(series[:8]):
        raw = str(sr.get("name", ""))
        color = HW if raw.split("·")[0].lower() in ("acme", "Acme") \
            else palette[i % len(palette)]
        # 图例带固定篮子件数（N 值恒定是口径承诺，让读者看得见）
        label = f"{raw}·{sr['n']}件" if sr.get("n") else raw
        pts = sr.get("pts") or []
        ax.plot(range(len(xs)), [p if p is not None else float("nan") for p in pts],
                marker="o", ms=4, lw=1.8, color=color, label=label)
        last = next(((k, v) for k, v in reversed(list(enumerate(pts)))
                     if v is not None), None)
        if last:
            ends.append((last[0], float(last[1]), color))
    # ★ 线末标价（2026-08-27 用户：「里面的均价都没有标注，没有价格」）。
    #   相邻标注按 y 间距排开，挤在一起等于没标
    if ends:
        span = (max(e[1] for e in ends) - min(e[1] for e in ends)) or 1.0
        min_gap = span * 0.055
        placed: list[float] = []
        for x0, y0, col in sorted(ends, key=lambda e: -e[1]):
            y_lab = y0
            while any(abs(y_lab - p) < min_gap for p in placed):
                y_lab -= min_gap
            placed.append(y_lab)
            ax.annotate(f"{y0:,.0f}", xy=(x0, y0), xytext=(x0 + 0.25, y_lab),
                        fontsize=8, color=col, fontweight="bold",
                        va="center", annotation_clip=False)
        ax.set_xlim(-0.5, len(xs) - 0.5 + max(1.6, len(xs) * 0.06))
    # 刻度太密一律抽稀：35 个日刻度挤成一团等于没有轴
    step = max(1, (len(xs) + 13) // 14)
    ax.set_xticks(range(0, len(xs), step))
    ax.set_xticklabels([xs[k] for k in range(0, len(xs), step)], fontsize=8.5)
    ax.set_ylabel(opt.get("ylab") or ylab or "", fontsize=8.5)
    ax.grid(axis="y", color="#E5E5EA", lw=0.7)
    ax.set_axisbelow(True)
    ax.legend(fontsize=8, ncol=min(4, len(series)), frameon=False,
              loc="lower center", bbox_to_anchor=(0.5, 1.0))
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    ax.tick_params(labelsize=8.5)
    return _fig_png(fig)


def _chart_share(opt: dict) -> bytes:
    """100% 堆叠：构成。与网页 share 同语义。"""
    rows = opt.get("rows") or []
    order = opt.get("order") or []
    if not rows or not order:
        return b""
    labels = [r.get("label", "") for r in rows]
    fig, ax = plt.subplots(figsize=(8, max(2.0, 0.5 * len(rows) + 1)))
    left = [0.0] * len(rows)
    palette = [C_DOWN, C_UP, C_ACCENT, C_GREY]
    for idx, seg in enumerate(order):
        vals, tot = [], []
        for r in rows:
            t = sum(float(r.get(s["k"]) or 0) for s in order) or 1.0
            tot.append(t)
            vals.append(float(r.get(seg["k"]) or 0) / t * 100)
        ax.barh(labels, vals, left=left, height=0.55,
                color=palette[idx % len(palette)], label=seg.get("name", ""))
        left = [l + v for l, v in zip(left, vals)]
    ax.set_xlim(0, 100)
    ax.set_xlabel("占比 %", fontsize=9)
    ax.legend(fontsize=8, ncol=len(order), frameon=False,
              loc="lower center", bbox_to_anchor=(0.5, 1.0))
    for s in ("top", "right", "left"):
        ax.spines[s].set_visible(False)
    ax.tick_params(labelsize=9)
    return _fig_png(fig)


def render_charts(charts: list[dict]) -> dict[str, bytes]:
    """把报告里的图渲成 PNG。画不出来的**跳过而不是塞占位图** ——
    导出物里出现一张空白图，比没有这张图更让人困惑。"""
    out: dict[str, bytes] = {}
    for c in charts or []:
        try:
            q, opt = c.get("question"), c.get("opt") or {}
            if q == "change":
                png = _chart_change(opt, c.get("xlab") or "")
            elif q == "deviation":
                png = _chart_deviation(opt, c.get("xlab") or "")
            elif q == "compare":
                png = _chart_compare(opt, c.get("xlab") or "")
            elif q == "share":
                png = _chart_share(opt)
            else:
                continue
            if png:
                out[c["el"]] = png
        except Exception:                              # noqa: BLE001
            continue
    return out


# ══════════════════════════ Word ══════════════════════════

def _docx_yahei(doc) -> None:
    """全文档统一微软雅黑（2026-08-27 用户点名字体不好看）。

    ★ 只设 style.font.name 是不够的：Word 对中文取的是 rFonts 的
      **w:eastAsia** 属性，不设它标题/表格仍会回落到宋体/等线 ——
      而且不报错，文档"看起来生成成功了"。收尾对**每一个 run**
      补 eastAsia，含表格单元格里的。
    """
    from docx.oxml.ns import qn

    def fix(runs):
        for r in runs:
            r.font.name = "Microsoft YaHei"
            r.element.get_or_add_rPr().get_or_add_rFonts().set(
                qn("w:eastAsia"), "微软雅黑")

    for p in doc.paragraphs:
        fix(p.runs)
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    fix(p.runs)


def to_docx(title: str, subtitle: str, md: str, charts: list[dict]) -> bytes:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    imgs = render_charts(charts)
    doc = Document()
    st = doc.styles["Normal"]
    st.font.name = "Microsoft YaHei"
    st.font.size = Pt(10.5)

    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.LEFT
    if subtitle:
        p = doc.add_paragraph(subtitle)
        p.runs[0].font.size = Pt(9)
        p.runs[0].font.color.rgb = RGBColor(0x6E, 0x6E, 0x73)

    for b in parse_blocks(md):
        if b["t"] == "h":
            doc.add_heading(b["text"], level=min(b["level"], 4))
        elif b["t"] == "p":
            doc.add_paragraph(b["text"])
        elif b["t"] == "quote":
            p = doc.add_paragraph(b["text"])
            p.runs[0].font.size = Pt(9)
            p.runs[0].font.color.rgb = RGBColor(0x6E, 0x6E, 0x73)
        elif b["t"] == "table":
            t = doc.add_table(rows=1, cols=len(b["head"]))
            t.style = "Light Grid Accent 1"
            for j, c in enumerate(b["head"]):
                cell = t.rows[0].cells[j]
                cell.text = c
                for r in cell.paragraphs[0].runs:
                    r.font.bold = True
                    r.font.size = Pt(9)
            for row in b["rows"]:
                cells = t.add_row().cells
                for j, c in enumerate(row[:len(b["head"])]):
                    cells[j].text = c
                    for r in cells[j].paragraphs[0].runs:
                        r.font.size = Pt(8.5)
            doc.add_paragraph()
        elif b["t"] == "chart":
            png = imgs.get(b["el"])
            if png:
                doc.add_picture(io.BytesIO(png), width=Inches(6.2))
        elif b["t"] == "hr":
            doc.add_paragraph("—" * 30)

    _docx_yahei(doc)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ══════════════════════════ PPT ══════════════════════════

def to_pptx(title: str, subtitle: str, md: str, charts: list[dict]) -> bytes:
    """每页回答一个问题，**标题就是结论**。

    版面铁律（2026-08-25 用户四条意见，逐条对应）：
    1. 页标题直接用报告里的结论式标题（#### 由 Agent 写好），
       不许出现「关键变化」「趋势与构成」这类看不出内容的通用标题；
    2. 图 + 表同页竖排，内容区排到 7.3 英寸 —— 不留半页空白；
    3. 每页右上角有「模块 · 国家」上下文条，标题里再带国家名，
       任何一页单独截出来都知道在讲哪国；
    4. 周线不足 4 周的问题在生成端解决（日线退化），这里只管画。
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image as PILImage

    imgs = render_charts(charts)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    INK = RGBColor(0x14, 0x23, 0x2B)
    GREY = RGBColor(0x6E, 0x6E, 0x73)
    RED = RGBColor(0xC7, 0x00, 0x0B)          # Acme身份色，只作装饰条
    FONT = _FONT or "Microsoft YaHei"

    def _txt(s, x, y, w, h, text, size, bold=False, color=INK, align=None,
             wrap=True):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.text = text
        for p in tf.paragraphs:
            if align is not None:
                p.alignment = align
            for r in p.runs:
                r.font.size, r.font.bold, r.font.name = Pt(size), bold, FONT
                r.font.color.rgb = color
        return tb

    def _rich(s, x, y, w, h, rows, gap=4):
        """一个文本框里多段不同样式的行。rows: (text, size, bold, color)"""
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        for k, (t, size, bold, color) in enumerate(rows):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.text = t
            p.space_after = Pt(gap)
            for r in p.runs:
                r.font.size, r.font.bold, r.font.name = Pt(size), bold, FONT
                r.font.color.rgb = color
        return tb

    def _fit_size(text, width_in, max_size=22, min_size=12):
        """算出让标题**正好一行放下**的字号。

        ★ 不能用 len() 阶梯：字号台阶粗、CJK 和「品牌名+百分号」宽度差一倍，
          50 个字落在 15pt 档照样折行 —— 用户点名「不能一次写清楚吗」。
          逐字符估宽（CJK/全角≈1 个字号宽，拉丁/数字≈0.55），
          直接反解字号，宁可字小一号也不许折行。
        """
        w = 0.0
        for ch in text:
            w += 1.0 if ord(ch) > 0x2E7F else 0.55
        pts = width_in * 72 * 0.96          # 4% 余量抵消估宽误差
        return max(min_size, min(max_size, int(pts / max(w, 1.0))))

    def page(headline, ctx=""):
        s = prs.slides.add_slide(blank)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 prs.slide_width, Inches(0.07))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RED
        bar.line.fill.background()
        bar.shadow.inherit = False
        if ctx:
            _txt(s, 10.83, 0.16, 2.0, 0.4, ctx, 11, color=GREY,
                 align=PP_ALIGN.RIGHT)
        tw = 10.2 if ctx else 12.3
        _txt(s, 0.5, 0.22, tw, 0.9, headline, _fit_size(headline, tw),
             bold=True, wrap=False)
        return s

    def _cell(cell, text, size, bold=False):
        cell.text = str(text)
        cell.margin_top = cell.margin_bottom = Pt(2)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size, r.font.bold, r.font.name = Pt(size), bold, FONT

    def put_table(s, head, rows, top, bottom=7.3):
        """列宽按内容长度加权（机型/价格列宽、日期列窄），行数按剩余空间截。"""
        avail = max(1, int((bottom - top - 0.34) / 0.31))
        cut = rows[:avail]
        ncol = len(head)
        shape = s.shapes.add_table(len(cut) + 1, ncol, Inches(0.5), Inches(top),
                                   Inches(12.3), Inches(0.31 * (len(cut) + 1)))
        tbl = shape.table
        wts = []
        for j in range(ncol):
            longest = max([len(str(head[j]))]
                          + [len(str(r[j])) for r in cut if j < len(r)])
            wts.append(min(longest, 28) + 2)
        tot = sum(wts) or 1
        for j in range(ncol):
            tbl.columns[j].width = Inches(12.3 * wts[j] / tot)
        for j, c in enumerate(head):
            _cell(tbl.cell(0, j), c, 11, bold=True)
        for i2, row in enumerate(cut, start=1):
            for j in range(ncol):
                _cell(tbl.cell(i2, j), row[j] if j < len(row) else "", 10.5)
        y_end = top + 0.31 * (len(cut) + 1)
        if len(rows) > len(cut):
            _txt(s, 0.5, min(y_end + 0.02, 7.05), 12.3, 0.3,
                 f"共 {len(rows)} 行，按重要度取前 {len(cut)} 行（全量见网页/PDF 版）",
                 9.5, color=GREY)
        return y_end

    def put_chart(s, png, top, bottom):
        w0, h0 = PILImage.open(io.BytesIO(png)).size
        h = min(bottom - top, 12.3 * h0 / w0)
        w = h * w0 / h0
        s.shapes.add_picture(io.BytesIO(png), Inches((13.333 - w) / 2),
                             Inches(top), height=Inches(h))
        return top + h

    def sentences(text):
        return [t.strip() for t in re.split(r"(?<=[。！？])", text or "")
                if t.strip()]

    blocks = parse_blocks(md)
    foot: list[str] = []      # 口径类内容统一收到最后一页，不挤占内容页

    # ── 封面 ──
    s = prs.slides.add_slide(blank)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                              prs.slide_width, Inches(0.16))
    band.fill.solid()
    band.fill.fore_color.rgb = RED
    band.line.fill.background()
    band.shadow.inherit = False
    _txt(s, 0.6, 2.35, 12.1, 1.3, title, 34, bold=True)
    if subtitle:
        _txt(s, 0.6, 3.7, 12.1, 0.6, subtitle, 14, color=GREY)

    # ── 本期结论：第一个 ## 之前的段落，按句拆成要点 ──
    lead: list[str] = []
    first_h2 = len(blocks)
    for idx, b in enumerate(blocks):
        if b["t"] == "h" and b["level"] == 2:
            first_h2 = idx
            break
        if b["t"] == "p":
            lead.append(b["text"])
    if lead:
        pts = []
        for t in lead:
            pts.extend(sentences(t))
        s = page("本期结论", "摘要")
        _rich(s, 0.6, 1.25, 12.1, 5.95,
              [("• " + t, 14, False, INK) for t in pts[:9]], gap=8)

    def h3_country(text):
        m = re.match(r"^(.+?)（([A-Z]{2})）\s*$", text or "")
        return (m.group(1), m.group(2)) if m else (text, "")

    # ── 主体：按 ## 模块分段处理 ──
    i, n = first_h2, len(blocks)
    while i < n:
        b = blocks[i]
        if b["t"] == "hr":
            i += 1
            continue
        if b["t"] == "quote":
            if b["text"].startswith("口径") or b["text"].startswith("ASP 图"):
                foot.append(b["text"])
            i += 1
            continue
        if not (b["t"] == "h" and b["level"] == 2):
            # 落在 hr 之后的收尾段落也进口径页
            if b["t"] == "p":
                foot.append(b["text"])
            i += 1
            continue

        h2 = b["text"]
        j = i + 1
        while j < n and not (blocks[j]["t"] == "h" and blocks[j]["level"] == 2):
            j += 1
        body = blocks[i + 1:j]

        # hr 之后是全文脚注区（最后一个模块的 body 会把它带进来，
        # 最后一个模块是谁都可能）—— 必须在分发进任何分支**之前**切走，
        # 否则口径文字会混进该分支当正文（VOC 页实测中过招）
        if any(x["t"] == "hr" for x in body):
            cut = next(k for k, x in enumerate(body) if x["t"] == "hr")
            for x in body[cut:]:
                if x["t"] in ("p", "quote"):
                    foot.append(x["text"])
            body = body[:cut]

        # ★ VOC 判定必须在「分国」之前：VOC 模块标题是「分国 VOC…」，
        #   先查"分国"会把口碑页整段吞进竞争分支（实测 5 页 VOC 全走错）
        if "VOC" in h2.upper():
            _pptx_voc_section(body, page, put_table, _rich, foot,
                              h3_country, INK, GREY)
        elif "分国" in h2:
            _pptx_country_section(body, page, put_table, put_chart, _rich,
                                  imgs, foot, h3_country, INK, GREY)
        else:
            # 通用模块：段落要点 + 表/图，同页竖排
            paras = [x["text"] for x in body if x["t"] == "p"]
            tables = [x for x in body if x["t"] == "table"]
            charts_ = [x for x in body if x["t"] == "chart" and imgs.get(x["el"])]
            for x in body:
                if x["t"] == "quote":
                    (foot if x["text"].startswith("口径") else paras).append(x["text"])
            if not (paras or tables or charts_):
                i = j
                continue
            s = page(h2)
            y = 1.2
            if paras:
                k = min(3, len(paras))
                _rich(s, 0.6, y, 12.1, 0.42 * k,
                      [("• " + t, 12.5, False, INK) for t in paras[:k]], gap=4)
                y += 0.42 * k + 0.1
            if charts_:
                y = put_chart(s, imgs[charts_[0]["el"]], y,
                              7.3 if not tables else 5.2) + 0.12
            if tables:
                put_table(s, tables[0]["head"], tables[0]["rows"], y)
            # 一个模块塞不下的第二张表/图，各自成页，标题带（续）
            for t2 in tables[1:]:
                s2 = page(h2 + "（续）")
                put_table(s2, t2["head"], t2["rows"], 1.25)
            for c2 in charts_[1:]:
                s2 = page(h2 + "（续）")
                put_chart(s2, imgs[c2["el"]], 1.25, 7.3)
        i = j

    # ── 口径与数据说明：全部脚注一页收尾 ──
    if foot:
        s = page("口径与数据说明", "附注")
        _rich(s, 0.6, 1.25, 12.1, 6.0,
              [("• " + t, 11.5, False, GREY) for t in foot[:10]], gap=6)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pptx_country_section(body, page, put_table, put_chart, _rich,
                          imgs, foot, h3_country, INK, GREY):
    """分组模块 → 每组一张「要点页」+ 每个小节一张「图+表页」。

    2026-08-27 起纲目是 产业(###) → 国家(####)；此前是 国家 → 品类。
    两种都认：h3 带「（XX）」国家码就按国家组，否则按产业组。
    页标题 = 组名 · 小节结论（#### 原文），右上角 = 模块 · 组 —— 单页可定位。
    """
    i, n = 0, len(body)
    cname, cc = "", ""
    ana = ""
    cat_heads: list[str] = []       # 该国各品类结论，攒给要点页
    pend: list[tuple] = []          # (标题, chart_el|None, table|None)

    def flush_country():
        nonlocal ana, cat_heads, pend
        if not (ana or pend):
            return
        # 国家要点页：分析段 + 各品类结论清单（也是该国的目录）。
        # 页标题用分析段第一句 —— 它就是该国最要紧的结论；
        # 已当标题的那句不再进正文，标题里也不重复国家名
        segs = [x.strip() for x in re.split(r"(?<=[。；])", ana) if x.strip()]
        lead1 = (segs[0].rstrip("。；") if segs else "").strip()
        ctx = f"分国竞争 · {cc}" if cc else f"分产业竞争 · {cname}"
        if 0 < len(lead1) <= 48:
            head_txt = lead1 if lead1.startswith(cname) else f"{cname}：{lead1}"
            body_segs = segs[1:5]
        else:
            head_txt = f"{cname}：近 30 天竞争要点"
            body_segs = segs[:4]
        s = page(head_txt, ctx)
        rows = [(t, 15, False, INK) for t in body_segs]
        rows += [("▸ " + h, 14, True, INK) for h in cat_heads[:6]]
        _rich(s, 0.6, 1.5, 12.1, 5.7, rows, gap=13)
        # 小节页：ASP 图 + 表同页；单品走势图独立成页（三样同页必然挤糊）
        for head, els, tbl in pend:
            s2 = page(f"{cname}·{head}", ctx)
            y = 1.2
            if els and imgs.get(els[0]):
                y = put_chart(s2, imgs[els[0]], y, 5.25 if tbl else 7.3) + 0.12
            if tbl:
                put_table(s2, tbl["head"], tbl["rows"], y)
            for el2 in els[1:]:
                if imgs.get(el2):
                    who = head.split(" ｜ ")[0]
                    s3 = page(f"{cname}·{who} ｜ 重点单品走势（台阶=真调价，脉冲=短促）",
                              ctx)
                    put_chart(s3, imgs[el2], 1.25, 7.3)
        ana, cat_heads, pend = "", [], []

    while i < n:
        b = body[i]
        if b["t"] == "h" and b["level"] == 3:
            flush_country()
            cname, cc = h3_country(b["text"])
            i += 1
            # 国家标题后的第一段是分析
            if i < n and body[i]["t"] == "p":
                ana = body[i]["text"]
                i += 1
            continue
        if b["t"] == "h" and b["level"] == 4:
            head = b["text"]
            cat_heads.append(head)
            els, tbl = [], None
            i += 1
            # p 也要跨过去：单品走势图前有一行加粗引导语，不跨会把图漏掉
            while i < n and body[i]["t"] in ("chart", "table", "p"):
                if body[i]["t"] == "chart":
                    els.append(body[i]["el"])
                elif body[i]["t"] == "table":
                    tbl = body[i]
                i += 1
            pend.append((head, els, tbl))
            continue
        if b["t"] == "quote":
            foot.append(b["text"])
        i += 1
    flush_country()


def _pptx_voc_section(body, page, put_table, _rich, foot, h3_country, INK, GREY):
    """VOC 模块 → 每国「统计页」+「原声页」。

    2026-08-27 用户要求 VOC 要统计不要罗列，于是这一段多了分产业统计表与
    品牌口碑表 —— 这里必须渲染表格（此前只认段落与引文，表被整个丢掉）。
    """
    i, n = 0, len(body)
    cname, cc = "", ""
    rows: list[tuple] = []
    tables: list[dict] = []
    stats = ""

    def flush():
        nonlocal rows, tables, stats
        if not (rows or tables):
            return
        head = f"VOC · {cname}：{stats}" if stats else f"VOC · {cname}"
        # ① 统计页：分产业表 + 品牌表（两张竖排，装不下的自动截行并注明）
        if tables:
            s = page(head, f"消费者声音 · {cc}")
            y = 1.25
            for t in tables[:2]:
                bottom = 4.2 if len(tables) > 1 and t is tables[0] else 7.3
                y = put_table(s, t["head"], t["rows"], y, bottom) + 0.28
        # ② 原声页：引文保持原文，译文灰字缩排
        if rows:
            s2 = page(f"{head} ｜ 消费者原声" if tables else head,
                      f"消费者声音 · {cc}")
            _rich(s2, 0.6, 1.4, 12.1, 5.9, rows[:14], gap=6)
        rows, tables, stats = [], [], ""

    while i < n:
        b = body[i]
        if b["t"] == "h" and b["level"] == 3:
            flush()
            cname, cc = h3_country(b["text"])
        elif b["t"] == "table":
            tables.append(b)
        elif b["t"] == "p":
            t = b["text"]
            if t.startswith("本期"):
                stats = t.rstrip("。")
            else:
                rows.append((t, 13.5, True, INK))
        elif b["t"] == "quote":
            t = b["text"]
            if t.startswith("口径"):
                foot.append(t)
            elif t.startswith("译："):
                rows.append(("　" + t, 12, False, GREY))
            else:
                rows.append((t, 12.5, False, INK))
        i += 1
    flush()


# ══════════════════════════ PDF ══════════════════════════

def to_pdf(title: str, subtitle: str, md: str, charts: list[dict]) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import (Image, Paragraph, SimpleDocTemplate,
                                    Spacer, Table, TableStyle)

    # ★ 中文必须显式注册字体，否则 reportlab 输出的中文是空白 —— 且不报错。
    # 首选微软雅黑 TTC（2026-08-27 用户点名 STSong 宋体不好看）；
    # 注册 bold 族让 <b> 标签真加粗而不是描边。雅黑不可用再退 STSong。
    cn = cnb = "Helvetica"
    try:
        from reportlab.pdfbase.ttfonts import TTFont
        pdfmetrics.registerFont(TTFont("MSYaHei", r"C:\Windows\Fonts\msyh.ttc",
                                       subfontIndex=0))
        cn = "MSYaHei"
        try:
            pdfmetrics.registerFont(TTFont("MSYaHei-Bold",
                                           r"C:\Windows\Fonts\msyhbd.ttc",
                                           subfontIndex=0))
            cnb = "MSYaHei-Bold"
        except Exception:                               # noqa: BLE001
            cnb = "MSYaHei"
        pdfmetrics.registerFontFamily("MSYaHei", normal=cn, bold=cnb,
                                      italic=cn, boldItalic=cnb)
    except Exception:                                   # noqa: BLE001
        try:
            pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
            cn = cnb = "STSong-Light"
        except Exception:                               # noqa: BLE001
            pass

    imgs = render_charts(charts)
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=18 * mm, rightMargin=18 * mm,
                            topMargin=16 * mm, bottomMargin=16 * mm,
                            title=title)
    ss = getSampleStyleSheet()
    S = {
        "title": ParagraphStyle("t", parent=ss["Title"], fontName=cnb,
                                fontSize=18, leading=24, alignment=0),
        "sub": ParagraphStyle("s", parent=ss["Normal"], fontName=cn,
                              fontSize=8.5, textColor=colors.HexColor("#6E6E73")),
        "h": ParagraphStyle("h", parent=ss["Heading2"], fontName=cnb,
                            fontSize=12.5, leading=17, spaceBefore=9, spaceAfter=4),
        "h3": ParagraphStyle("h3", parent=ss["Heading3"], fontName=cnb,
                             fontSize=11, leading=15, spaceBefore=7, spaceAfter=3),
        "p": ParagraphStyle("p", parent=ss["Normal"], fontName=cn,
                            fontSize=9.5, leading=15),
        "cell": ParagraphStyle("c", parent=ss["Normal"], fontName=cn,
                               fontSize=7.6, leading=10),
    }
    flow = [Paragraph(title, S["title"])]
    if subtitle:
        flow += [Spacer(1, 3), Paragraph(subtitle, S["sub"])]
    flow.append(Spacer(1, 8))

    avail = A4[0] - 36 * mm
    for b in parse_blocks(md):
        if b["t"] == "h":
            flow.append(Paragraph(b["text"], S["h" if b["level"] <= 2 else "h3"]))
        elif b["t"] in ("p", "quote"):
            flow.append(Paragraph(b["text"], S["p" if b["t"] == "p" else "sub"]))
            flow.append(Spacer(1, 3))
        elif b["t"] == "table":
            ncol = len(b["head"])
            data = [[Paragraph(f"<b>{c}</b>", S["cell"]) for c in b["head"]]]
            for row in b["rows"]:
                data.append([Paragraph(c, S["cell"]) for c in row[:ncol]])
            t = Table(data, colWidths=[avail / ncol] * ncol, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#F2F2F7")),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#D1D1D6")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]))
            flow += [t, Spacer(1, 8)]
        elif b["t"] == "chart":
            png = imgs.get(b["el"])
            if png:
                from PIL import Image as PILImage
                w, h = PILImage.open(io.BytesIO(png)).size
                iw = avail
                flow += [Image(io.BytesIO(png), width=iw, height=iw * h / w),
                         Spacer(1, 8)]
    doc.build(flow)
    return buf.getvalue()


EXPORTERS = {"docx": to_docx, "pptx": to_pptx, "pdf": to_pdf}
MIME = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
}


def export(fmt: str, title: str, subtitle: str, md: str,
           charts: list[dict]) -> tuple[bytes, str]:
    fn = EXPORTERS.get((fmt or "").lower())
    if not fn:
        raise ValueError(f"不支持的格式：{fmt}（可用：{'/'.join(EXPORTERS)}）")
    data = fn(title, subtitle, md, charts)
    safe = re.sub(r"[^\w一-龥-]+", "_", title)[:60] or "report"
    stamp = datetime.now().strftime("%Y%m%d")
    return data, f"{safe}_{stamp}.{fmt.lower()}"
