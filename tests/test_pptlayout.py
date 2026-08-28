# -*- coding: utf-8 -*-
"""PPT 导出版式与价格曲线口径的回归测试。

用户的意见（2026-08-25，原话拆四条）：
  「每页我都不知道你要表达什么，让Agent写标题啊」
  「每一页的结构也太不紧凑了，留白留那么多干啥」
  「每页我都不知道是哪个国家的」
  「周度的价格，你得连续几周，至少4周才能看出来吧」

同轮实测踩的坑（都不报错、图都画得出来，只有断言守得住）：
  - 日线 ASP 直接 AVG 当日观测 → 构成效应「+590%」假涨价（第 5 次犯，
    见 knowledge/lessons/composition-vs-level-effects.md）
  - 基期准入按全窗口卡 → 把晚入库的Acme整条踢出自家图
  - LLM 标题长度闸 40 对混排太紧 → 合格标题全被误杀回兜底
  - 「分国 VOC」标题先匹配到"分国"分支 → 5 页 VOC 全走错模块
"""
import io
import os
import pathlib
import re
import sys

ROOT = pathlib.Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
os.environ.setdefault("PYTHONUTF8", "1")

FAIL, PASS = [], [0]


def ok(cond, msg):
    if cond:
        PASS[0] += 1
    else:
        FAIL.append(msg)


# ─────────── 1. 功能测试：拼一份迷你报告，验证 PPT 结构 ───────────
from app import report_export  # noqa: E402

MD = """# 测试双周报

> 双周报 · 2026-08-20 ~ 2026-09-04 · 5 国

本期 Samsung 在智利领跌。Acme均价守稳。

## ⚠ 价格预警：1 起打到我方对位，最大 Redmi 15C -29%（Bodega）

| 机型 | 国家 | 变动 |
|---|---|---|
| Redmi 15C | MX | -29% |

## 分产业 × 分国家竞争分析（近 30 天）

### 手机
手机品类竞争激烈。变动集中在 Falabella。

#### 智利 ｜ 三星领跌，均价普降3%-5%
![chart:prod-CL-phone]

| 机型 | 变动 |
|---|---|
| Galaxy A07 | +15% |

## 🗣 分国 VOC（新品口碑 · 含原声）

### 智利（CL）
本期收集 10 条新品评论（差评 4，40%）。

**Galaxy A17**★新 · 👎 差评 · Falabella
> 「Necesito ayuda para traspasar información」（西语）
> 译：需要帮助转移数据

> 口径：近 30 天收集。

---
口径：正文只讲重要度最高的 6 条。
"""

CHARTS = [{"question": "change", "el": "prod-CL-phone",
           "title": "智利·手机：三星领跌",
           "opt": {"xs": ["08-10", "08-11", "08-12"],
                   "series": [{"name": "Samsung", "pts": [700, 690, 680]},
                              {"name": "Acme", "pts": [660, 662, 665]}],
                   "ylab": "挂牌价 CLP"}}]

data = report_export.to_pptx("测试双周报", "2026-08-20 ~ 2026-09-04", MD, CHARTS)
ok(len(data) > 10000, "pptx 要能生成")

from pptx import Presentation  # noqa: E402

prs = Presentation(io.BytesIO(data))
slide_texts = []
for s in prs.slides:
    txts = [sh.text_frame.text for sh in s.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()]
    pics = sum(1 for sh in s.shapes if sh.shape_type == 13)
    tbls = sum(1 for sh in s.shapes if sh.has_table)
    slide_texts.append({"txts": txts, "pics": pics, "tbls": tbls,
                        "all": " ".join(txts)})

ok(all(st["txts"] for st in slide_texts), "★ 每一页都必须有文字（不许空板）")
ok(not any("关键变化" in st["all"] or "趋势与构成" in st["all"]
           for st in slide_texts),
   "★ 通用标题「关键变化 / 趋势与构成」必须绝迹 —— 用户：每页都不知道要表达什么")

# 国家要点页的目录 bullet 里也有这句结论，必须挑**带图**的那页
cat = next((st for st in slide_texts
            if "三星领跌" in st["all"] and st["pics"] >= 1), None)
ok(cat is not None, "品类页要存在（标题含 Agent 结论）")
if cat:
    ok(cat["pics"] == 1 and cat["tbls"] == 1,
       "★ 小节页图和表必须同页竖排 —— 不许一图一页留半页空白")
    ok("智利" in cat["all"] and "手机" in cat["all"],
       "★ 小节页标题必须同时带产业与国家（2026-08-27 起产业为纲）")
    ok("分产业竞争 · 手机" in cat["all"],
       "★ 每页右上角要有「模块 · 组」上下文")

voc = next((st for st in slide_texts if "VOC" in st["all"]), None)
ok(voc is not None, "VOC 页要存在")
if voc:
    ok("消费者声音" in voc["all"],
       "★ VOC 页要走 VOC 分支 —— 「分国 VOC」曾被'分国'先匹配走错模块")
    ok("Necesito ayuda" in voc["all"], "VOC 页要有西语原声")
    ok("译：" in voc["all"], "VOC 页要有译文")

alert = next((st for st in slide_texts if "价格预警" in st["all"]), None)
ok(alert is not None and alert["tbls"] == 1, "预警页标题即结论 + 带表")

foot_sl = next((st for st in slide_texts if "口径与数据说明" in st["all"]), None)
ok(foot_sl is not None, "口径类内容要归到最后的附注页")
ok(not any("口径：正文只讲" in st["all"] for st in slide_texts
           if "口径与数据说明" not in st["all"]),
   "★ hr 之后的全文脚注不许漏进内容页当要点（市场动态页曾中招）")

# ─────────── 1b. 标题必须一行放下（用户：为什么要换行） ───────────
rx = (ROOT / "app/report_export.py").read_text(encoding="utf-8")
ok("_fit_size" in rx and "wrap=False" in rx,
   "★ 页标题按字宽反解字号 + 禁止折行 —— len() 阶梯档太粗，"
   "50 字混排在 15pt 档照样折成两行（用户截图点名）")
ok("0x2E7F" in rx,
   "字宽估算要区分 CJK/全角（≈1）与拉丁/数字（≈0.55），"
   "「品牌名+百分号」按 CJK 算会把字号压得没必要地小")

# ─────────── 2. 源码断言：产品价格曲线（品牌均价已废弃） ───────────
wk = (ROOT / "app/agents/weekly.py").read_text(encoding="utf-8")
ok("_asp_weekly_chart" not in wk,
   "★ 品牌均价图必须彻底删除 —— 2026-08-27 用户：「所有的友商均价肯定都"
   "不太准，你就举例看具体产品的价格就行了」。实测墨西哥三星手机 134 个 "
   "SKU 均值 624 美元、中位数仅 376，被折叠屏的多个变体拉飞；"
   "无销量权重时它衡量的是货架构成而非成交价")
ok("fx_usd" not in wk and "def _fx" not in wk,
   "均价没了，配套的美元换算也不该留在周报里（跨国比较改看具体产品）")
m = re.search(r"def _product_trend_chart.*?\n    def ", wk, re.S)
ok(m is not None, "找不到 _product_trend_chart")
src = m.group(0) if m else ""
ok("_top_products" in src,
   "★ 达标变动不足时要用主力在售产品补位 —— 否则没有变动的小节整格空白")
ok("_PROD_LINES" in src, "每张图的产品条数要有统一上限")
ok('"cur"' in src and "混币种" in src, "一张图一个币种（本币计价，不做换算）")
ok("html" in src and "unescape" in src,
   "商品标题里的 HTML 实体要解码，否则图例出现 'ASUS 14&quot'")

mtp = re.search(r"def _top_products.*?\n    def ", wk, re.S)
srctp = mtp.group(0) if mtp else ""
ok("dz >= 5" in srctp, "补位产品要观测天数够（线才画得出来）")
ok("STALE_DAYS" in srctp, "补位产品要排除老品")
ok("seen_model" in srctp, "同一机型只取一个渠道，避免一张图全是同一台机器")

m2 = re.search(r"def _cat_headlines.*?\n    def ", wk, re.S)
src2 = m2.group(0) if m2 else ""
ok(m2 is not None and "_ask_prose" in src2,
   "★ 小节标题要让 Agent 写（一次批量调用），不是程序拼")
ok("order_keys" in src2,
   "★ 模型抄字面 KEY 时按行序兜底（实测发生过）")
ok("64" in src2,
   "★ 标题长度闸 ≥64：40 对「拉美品牌名+百分号」混排太紧，曾整批误杀")
ok("det[" in src2 or "det =" in src2, "每个 key 都要有确定性兜底")

# ─────────── 2b. 概览一致性与措辞（用户：「怎么可能没变动」） ───────────
ok("（在盯，没动静）" not in wk and "报告期内无新变动" in wk,
   "★ 交付物措辞必须书面 —— 「没动静」只许出现在给模型的禁令里，"
   "不许出现在任何输出字符串中")
m2b = re.search(r"def _brief_prose.*?\n    def ", wk, re.S)
src2b = m2b.group(0) if m2b else ""
ok("近 30 天累计" in src2b and "_moves30" in src2b,
   "★ 概览事实必须含近 30 天窗口 —— 只喂报告期会写出「市场无变动」，"
   "而三行之下就是 56 个变动的表")
ok("不得写成整个市场没有变动" in src2b,
   "报告期为 0 时只许说「暂无新增」，不许说市场没变动")
m2c = re.search(r"def _brief_focus.*?\n    def ", wk, re.S)
src2c = m2c.group(0) if m2c else ""
ok("last30" in src2c and "'-30 day'" in src2c,
   "★ 重点关注期内为空时必须回看 30 天补最近一次动作")
ok('COUNTRY_ORDER = ("MX", "CO", "CL", "PE", "AR", "BR")' in wk,
   "★ 分国顺序是用户指定的固定顺序，不按数据量排")
_ac = re.search(r"def _active_countries.*?\n    def ", wk, re.S)
ok(_ac and "_cc_sort" in _ac.group(0),
   "★ 国家全集要经固定顺序排序（分产业段走 _active_countries）")
ok("self._cc_sort(by_cc)" in wk,
   "VOC 模块也要走固定国家顺序")

# ─────────── 3. 报告结构：分产业为纲 + #### 结论标题进 markdown ───────────
m3 = re.search(r"def _brief_country_cat.*?\n    def ", wk, re.S)
src3 = m3.group(0) if m3 else ""
ok("分产业 × 分国家" in src3,
   "★ 纲目是 产业(###) → 国家(####)（2026-08-27 用户把顺序倒过来）")
ok('f"#### {names.get(cc, cc)} ｜ {head}"' in src3,
   "★ 国家小节是 ####，标题本身就是结论（网页/Word/PDF/PPT 四端共用）")
ok("不再提供品牌均价" in src3 and "延续上次价格" in src3,
   "★ 口径脚注要写明：图上是具体产品挂牌价 + LOCF 延续，"
   "且必须说清为什么不再提供品牌均价（无销量权重＝衡量货架构成）")
ok("_cat_analysis" in src3, "每个产业开头要有跨国分析段")

# 反混淆纪律：所有幅度都属于具体产品
m4 = re.search(r"def _cat_analysis.*?\n    def ", wk, re.S)
src4 = m4.group(0) if m4 else ""
ok("整体均价口径" in src4 or "不存在任何品牌或品类的整体均价" in src4,
   "★ 喂给模型的事实里不能有任何整体均价数字 —— 单品降幅曾被写成"
   "品类整体降幅（用户点名）")
ok("禁止出现" in src4 and "均价" in src4,
   "提示词要明令禁止「均价/整体降幅/普降」这类整体口径")
ok("在跟踪产品" in src4, "分析事实改喂图上在跟踪的具体产品名")
ok("绝不许" in src4 or "必须点名" in src4, "每个幅度都要点名机型与渠道")

# 价格变动检测必须在自动流水线里
orch = (ROOT / "app/agents/orchestrator.py").read_text(encoding="utf-8")
ok("PriceMoveAgent" in orch and "价格变动检测" in orch,
   "★ 价格变动检测必须是采集流水线的一个阶段 —— 它此前只能手动触发，"
   "结果 2026-08-13 之后 14 天没跑，price_obs 每天照常入库而 price_move "
   "一条没新增，周报核心静默停更且页面看不出异常")
ok(orch.index("价格变动检测") > orch.index("PriceAuditAgent("),
   "★ 变动检测要排在价格审计之后 —— 在未审计的脏数据上比价会播报"
   "根本不存在的降价")

appjs = (ROOT / "app/web/app.js").read_text(encoding="utf-8")
ok("#### " in appjs, "★ 前端 mdToHtml 要认 ####（否则网页显示井号原文）")

# ─────────── 4. 导出字体：微软雅黑（2026-08-27 用户点名） ───────────
ok("_docx_yahei" in rx and "w:eastAsia" in rx,
   "★ Word 必须逐 run 设 w:eastAsia=微软雅黑 —— 只设 style.font.name "
   "中文仍回落宋体/等线，且不报错")
ok("msyh.ttc" in rx and "subfontIndex" in rx,
   "★ PDF 首选微软雅黑 TTC（STSong 宋体被用户点名不好看），"
   "TTC 必须带 subfontIndex")
ok("registerFontFamily" in rx, "PDF 要注册 bold 族，<b> 才真加粗")
ok("STSong-Light" in rx, "雅黑不可用时要有 STSong 兜底，不能空白")

chjs = (ROOT / "app/web/charts.js").read_text(encoding="utf-8")
ok("disp(s)" in chjs and "cmap[s.name]" in chjs,
   "★ 网页图例带 n 件数但颜色按裸品牌名查 —— 否则Acme红等身份色全丢")

# ─────────── 5. 单品走势与价格标注（2026-08-27 用户两点意见） ───────────
ok("endLabel" in chjs and "moveOverlap" in chjs,
   "★ 网页折线要有线末价格标注 + 防重叠 —— 「图里没有价格，看不出结论」")
ok("annotate" in rx and "min_gap" in rx,
   "★ 导出折线同样要线末标价，且相邻标注要错开")
ok("_mover_series" in wk and "回升（短促）" in wk and "新价已维持" in wk,
   "★ 每个单品变动必须判定后续：维持N天=真调价台阶，回升=促销脉冲 —— "
   "「否则你怎么知道这个不是瞬间的值，而是一个长促」")
ok('"prod-' in wk or "prod-{cc}" in wk,
   "★ 讲了单品就要配单品历史曲线（prod- 图，本币计价，LOCF 延续）")
ok("变动后续" in wk, "Top 表要有「变动后续」列")
ok("台阶" in wk and "脉冲" in wk,
   "分析提示词要教模型区分台阶与脉冲，短促不当战略降价")

# ─────────── 6. 全矩阵 + VOC 统计（2026-08-27 用户第三轮意见） ───────────
ok("_active_countries" in wk and "_active_cats" in wk,
   "★ 国家/产业全集要从库里 enabled 取，不写死 —— 用户：「所有国家所有产业，"
   "而不是让你挑着放」")
ok("for c in cats:" in src3 and "for cc in ccs:" in src3,
   "★ 小节要遍历全矩阵（每个启用产业 × 每个启用国家），不是只遍历有变动的")
ok("_coverage30" in wk and "不是没有数据" in wk,
   "★ 空格子必须用观测底子说明是「价格平稳」而不是「没抓到」——"
   "这是本项目最贵的一类错")
ok("没有采到该国该品类的价格观测" in wk,
   "真的没观测时要明确标为采集缺口，不许写成友商没动作")
ok("采集缺口" in wk and "绝不能写成" in wk,
   "跨国分析提示词也要禁止把数据缺口说成没有动作")

voc_src = wk[wk.index("def _brief_voc_country"):wk.index("def _brief_country_cat")]
ok("review_aspect" in voc_src,
   "★ VOC 统计要基于 VOC Agent 的维度标注（review_aspect），不是关键词猜")
ok("_asp_rank" in voc_src and "主要抱怨" in voc_src,
   "★ VOC 要有统计（各维度抱怨条数），不是单纯罗列几条评论")
ok("for c in self._active_cats()" in voc_src,
   "★ VOC 必须按产业全集分组 —— 用户：「每个国家不能只放手机，或只放一个产业」")
ok("| 品牌 |" in voc_src or "品牌\", \"评论数" in voc_src,
   "VOC 要有品牌口碑对比表（谁在挨骂）")
ok("_is_real_voice" in voc_src and "estrellas" in voc_src,
   "★ 引用前必须过滤界面控件文本 —— 实测「5 estrellas 91 %」评分分布"
   "被当消费者原声引进过周报；引用错了比少引一条贵得多")
ok("per_cat_cap" in voc_src,
   "原声要跨品类挑（每品类限额），否则又变成只有手机")

print(f"pptlayout: {PASS[0]} 通过, {len(FAIL)} 失败")
for f in FAIL:
    print("  ✗", f)
sys.exit(1 if FAIL else 0)
